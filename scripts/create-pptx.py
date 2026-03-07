#!/usr/bin/env python3
"""달려라 꽃배달 - 주문 등록 개선 제안서 PPTX 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color palette ───
GREEN = RGBColor(0x10, 0xB9, 0x81)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color=WHITE, line_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_text(slide, items, left, top, w, h, size=13, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.space_after = Pt(4)
    return txBox

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, GREEN)

add_text(slide, '달려라 꽃배달', Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         size=20, color=RGBColor(0xD1,0xFA,0xE5))
add_text(slide, '주문 등록 페이지 개선 제안서', Inches(1), Inches(2.2), Inches(11), Inches(1.2),
         size=40, bold=True, color=WHITE)
add_text(slide, '경쟁사 분석 기반 UI/UX 리디자인', Inches(1), Inches(3.5), Inches(11), Inches(0.6),
         size=18, color=RGBColor(0xD1,0xFA,0xE5))

add_text(slide, '2026.03.07', Inches(1), Inches(5.5), Inches(3), Inches(0.5),
         size=14, color=RGBColor(0xA7,0xF3,0xD0))
add_text(slide, '비교 대상: 468.co.kr (꽃비 파트너스) / ebestflower.co.kr (베스트플라워)', 
         Inches(1), Inches(6.0), Inches(10), Inches(0.5),
         size=12, color=RGBColor(0xA7,0xF3,0xD0))

# ═══════════════════════════════════════════════════════════
# SLIDE 2: 현황 요약
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '1. 현재 현황', Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         size=28, bold=True, color=DARK)
add_text(slide, '현재 주문 등록 페이지는 단순 텍스트 파싱 스텁으로, 실제 주문 등록 기능이 없습니다.',
         Inches(0.8), Inches(1.0), Inches(11), Inches(0.5), size=14, color=GRAY)

# 3 cards side by side
labels = ['현재 (달려라 꽃배달)', '468.co.kr (꽃비)', 'ebestflower.co.kr (베스트)']
descs = [
    '• 텍스트 파싱 textarea 1개\n• 주문 등록 API 미연결\n• 필드 입력 폼 없음\n• 메시지 붙여넣기만 가능',
    '• PHP 기반 전통적 테이블 폼\n• 팝업 기반 검색/선택\n• 모든 필드 한 화면에 나열\n• 모바일 미지원\n• 최근 발주 불러오기 지원',
    '• PHP 기반 테이블 폼\n• 30+ 상품 카테고리\n• 옵션 상품 체크박스\n• 경조사어/보내는분 추가/삭제\n• 카드 메시지 24종 카테고리\n• 해피콜/현장사진 요청'
]
colors_top = [RED, AMBER, BLUE]

for i in range(3):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(1.7), Inches(3.8), Inches(4.8), WHITE, RGBColor(0xE2,0xE8,0xF0))
    # color bar top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(3.8), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors_top[i]
    bar.line.fill.background()
    
    add_text(slide, labels[i], x + Inches(0.2), Inches(2.15), Inches(3.4), Inches(0.4),
             size=14, bold=True, color=DARK)
    add_text(slide, descs[i], x + Inches(0.2), Inches(2.6), Inches(3.4), Inches(3.5),
             size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: 경쟁사 폼 필드 비교
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '2. 경쟁사 폼 필드 비교 분석', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)

# Table-like comparison
headers = ['구분', 'ebestflower (베스트)', '우리 구현 계획']
rows = [
    ['발주/수주 화원', '팝업 화원 선택, 자동입력', '관리자가 직접 등록 (불필요)'],
    ['상품 선택', '드롭다운 30개 + 상품검색 팝업', '칩(Chip) 버튼 18종 + 상세명'],
    ['금액', '원청액/결제액/배송비 분리 입력', '상품가 + 배송비 + 옵션 자동합산'],
    ['옵션', '체크박스 8종 + 개별 금액 입력', '체크박스 8종 + 개별 금액 입력'],
    ['주문자/받는분', '이름 + 전화 + 핸드폰 (각각)', '이름 + 전화 + 핸드폰 (2컬럼)'],
    ['배달일시', '년/월/일 셀렉트 + 시간 드롭다운', '날짜 피커 + 시간 드롭다운'],
    ['배달주소', '주소검색 팝업 + 병원/예식장 검색', 'Daum 주소 인라인 + 상세주소'],
    ['메시지 유형', '카드만/리본만/카드+리본 라디오', '동일 (라디오 Pill 버튼)'],
    ['경조사어', '검색팝업 + 추가/삭제 + 복사 버튼', '프리셋 버튼 10종 + 직접 입력'],
    ['카드 메시지', '24종 카테고리 + 검색 + textarea', '카테고리 선택 + textarea'],
    ['요구사항', '검색 + 현장사진 체크 + textarea', '체크박스 옵션 + textarea'],
    ['해피콜', '라디오(요청/안함) + 사진URL 체크', '체크박스 (해피콜/사진/URL)'],
    ['장례/예식 정보', '없음 (주소에 포함)', '조건부 표시 (근조→장례, 축하→예식)'],
]

# Draw table header
y = Inches(1.2)
col_widths = [Inches(2), Inches(4.8), Inches(4.8)]
col_x = [Inches(0.8), Inches(2.8), Inches(7.6)]

for i, h in enumerate(headers):
    shape = add_shape(slide, col_x[i], y, col_widths[i], Inches(0.4), GREEN)
    add_text(slide, h, col_x[i] + Inches(0.1), y + Inches(0.02), col_widths[i] - Inches(0.2), Inches(0.35),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    ry = y + Inches(0.4) + Inches(ri * 0.38)
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF8,0xFA,0xFC)
    for ci, cell in enumerate(row):
        shape = add_shape(slide, col_x[ci], ry, col_widths[ci], Inches(0.38), bg, RGBColor(0xE2,0xE8,0xF0))
        sz = 10 if ci == 0 else 10
        add_text(slide, cell, col_x[ci] + Inches(0.1), ry + Inches(0.02), col_widths[ci] - Inches(0.2), Inches(0.33),
                 size=sz, color=DARK if ci == 0 else GRAY, bold=(ci==0))

# ═══════════════════════════════════════════════════════════
# SLIDE 4: 개선 포인트
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '3. 핵심 개선 포인트 (10가지)', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)

improvements = [
    ('1', '카카오톡 메시지 파싱', '메시지 붙여넣기로 주문 필드 자동 채움\n(기존 기능 발전)', GREEN),
    ('2', '모던 카드 기반 UI', '테이블→카드 섹션 레이아웃\n큰 글씨(19px), 명확한 시각 구분', GREEN),
    ('3', '조건부 필드 표시', '근조 선택→장례식장/고인/상주 자동표시\n축하 선택→예식장/홀 자동표시', AMBER),
    ('4', 'Daum 주소 인라인 검색', '팝업 대신 인라인 주소 자동완성\n(경쟁사 대비 UX 우위)', BLUE),
    ('5', '실시간 금액 자동 계산', '상품가 + 옵션가 + 배송비\n총 결제액 실시간 표시', GREEN),
    ('6', '경조사어 프리셋 버튼', '자주 쓰는 10종 원클릭 선택\n직접 입력도 가능', GREEN),
    ('7', '전화번호 자동 포맷', '입력 시 자동 하이픈\n010-1234-5678 형식', GRAY),
    ('8', '최근 발주 불러오기', '이전 주문 기반 필드 자동 채움\n반복 입력 최소화', BLUE),
    ('9', 'Zod 실시간 유효성 검사', '필수 필드 표시 + 즉시 에러 피드백\n(react-hook-form + zod)', AMBER),
    ('10', '완전한 모바일 대응', '반응형 1컬럼 레이아웃\n(경쟁사는 모바일 미지원)', RED),
]

for i, (num, title, desc, clr) in enumerate(improvements):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    yy = Inches(1.2 + row * 1.18)
    
    card = add_shape(slide, x, yy, Inches(5.9), Inches(1.05), WHITE, RGBColor(0xE2,0xE8,0xF0))
    
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), yy + Inches(0.2), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = clr
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    
    add_text(slide, title, x + Inches(0.8), yy + Inches(0.08), Inches(4.8), Inches(0.35),
             size=14, bold=True, color=DARK)
    add_text(slide, desc, x + Inches(0.8), yy + Inches(0.42), Inches(4.8), Inches(0.55),
             size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: 화면 구조 (Wire‑frame)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '4. 화면 구조 설계', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)
add_text(slide, '7개 섹션 카드 기반, 순차적 입력 흐름', Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
         size=14, color=GRAY)

sections_info = [
    ('1', '상품 선택', '칩 버튼 18종\n상세 상품명 / 수량\n상품 사진 URL 또는 업로드', GREEN),
    ('2', '주문자 (보내는 분)', '이름\n전화번호\n핸드폰 (자동 하이픈)', BLUE),
    ('3', '받는 분', '이름 *\n전화번호 *\n핸드폰', BLUE),
    ('4', '배달 정보', '배달일 (Date Picker)\n배달 시간 (드롭다운)\n주소 검색 + 상세주소\n[조건부] 장례/예식 정보', AMBER),
    ('5', '리본 / 카드 메시지', '메시지 유형 선택\n경조사어 프리셋 + 직접입력\n리본 좌/우측\n카드 메시지 textarea', GREEN),
    ('6', '금액', '상품가 / 배송비 / 옵션\n옵션 체크박스 8종\n총 결제액 자동 합산', RED),
    ('7', '추가 정보', '해피콜/사진 체크박스\n요구사항 textarea', GRAY),
]

for i, (num, title, desc, clr) in enumerate(sections_info):
    if i < 4:
        x = Inches(0.8 + i * 3.1)
        yy = Inches(1.5)
    else:
        x = Inches(0.8 + (i - 4) * 3.9)
        yy = Inches(4.1)
    
    w = Inches(2.8) if i < 4 else Inches(3.6)
    h = Inches(2.3) if i < 4 else Inches(2.3)
    
    card = add_shape(slide, x, yy, w, h, WHITE, RGBColor(0xE2,0xE8,0xF0))
    
    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, yy, w, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    
    add_text(slide, f'{num}  {title}', x + Inches(0.15), yy + Inches(0.03), w - Inches(0.3), Inches(0.35),
             size=13, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.15), yy + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
             size=11, color=GRAY)

# Arrow flow
add_text(slide, '▶ 메시지 파싱 바 (상단 고정) — 카카오톡/SMS 붙여넣기로 모든 섹션 자동 채움',
         Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
         size=13, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: 상품 카테고리 & 옵션
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '5. 상품 카테고리 & 옵션 상품', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)

# Product categories table
categories = [
    ('축하 화환', ['축하 3단', '축하 4단', '축하 5단', '실속축하'], AMBER),
    ('근조 화환', ['근조 3단', '근조 4단', '근조 5단', '실속근조'], RGBColor(0x64,0x74,0x8B)),
    ('꽃/오브제', ['꽃바구니', '꽃다발', '꽃상자', '오브제'], RGBColor(0xEC,0x48,0x99)),
    ('란/식물', ['동양란', '서양란', '관엽식물', '-'], GREEN),
    ('기타', ['쌀화환', '과일바구니', '기타', '-'], BLUE),
]

for i, (cat, items, clr) in enumerate(categories):
    yy = Inches(1.3 + i * 0.7)
    
    bar = add_shape(slide, Inches(0.8), yy, Inches(2), Inches(0.55), clr)
    add_text(slide, cat, Inches(0.9), yy + Inches(0.07), Inches(1.8), Inches(0.4),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    for j, item in enumerate(items):
        if item == '-':
            continue
        ix = Inches(3.0 + j * 1.8)
        chip = add_shape(slide, ix, yy + Inches(0.05), Inches(1.6), Inches(0.42), WHITE, clr)
        add_text(slide, item, ix + Inches(0.1), yy + Inches(0.07), Inches(1.4), Inches(0.35),
                 size=11, color=DARK, align=PP_ALIGN.CENTER)

# Options
add_text(slide, '옵션 상품 (체크박스 선택 + 개별 금액 입력)', Inches(0.8), Inches(5.0), Inches(10), Inches(0.4),
         size=16, bold=True, color=DARK)

options = ['케익', '샴페인', '사탕', '초콜렛', '떡케익', '와인', '빼빼로', '기타']
for i, opt in enumerate(options):
    x = Inches(0.8 + i * 1.5)
    chip = add_shape(slide, x, Inches(5.5), Inches(1.3), Inches(0.5), WHITE, RGBColor(0xE2,0xE8,0xF0))
    add_text(slide, f'☐ {opt}', x + Inches(0.1), Inches(5.53), Inches(1.1), Inches(0.4),
             size=12, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: 경조사어 & 카드 메시지
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '6. 경조사어 프리셋 & 카드 메시지 카테고리', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)

# Preset buttons
add_text(slide, '경조사어 프리셋 (원클릭 선택)', Inches(0.8), Inches(1.2), Inches(6), Inches(0.4),
         size=16, bold=True, color=DARK)

presets = [
    '축 결혼', '축 개업', '축 승진', '축 취임', '축 당선',
    '축하합니다', '삼가 고인의 명복을 빕니다', '사랑합니다',
    '생일 축하합니다', '쾌유를 빕니다', '축 이사', '축 출산',
    '입학을 축하합니다', '졸업을 축하합니다', '어버이날 감사합니다',
]

for i, p in enumerate(presets):
    col = i % 5
    row = i // 5
    x = Inches(0.8 + col * 2.4)
    yy = Inches(1.7 + row * 0.5)
    w = Inches(2.2)
    chip = add_shape(slide, x, yy, w, Inches(0.38), WHITE, GREEN)
    add_text(slide, p, x + Inches(0.05), yy + Inches(0.02), w - Inches(0.1), Inches(0.33),
             size=11, color=DARK, align=PP_ALIGN.CENTER)

# Card message categories
add_text(slide, '카드 메시지 카테고리 (ebestflower 24종 기반)', Inches(0.8), Inches(3.6), Inches(8), Inches(0.4),
         size=16, bold=True, color=DARK)

card_cats = [
    '사랑의 메시지', '생일축하', '결혼기념일', '승진/영전', '취임/퇴임', '개업/창립',
    '이사축하', '전시/모임', '수상/우승', '선거당선', '수연/고희', '병문안/위로',
    '조문/애도', '크리스마스', '어버이/스승의날', '연하/명절', '돌/백일', '어린이날',
    '합격/축하', '입학/졸업', '성년의날', '군입/전역', '약혼축하', '출산축하',
]

for i, c in enumerate(card_cats):
    col = i % 6
    row = i // 6
    x = Inches(0.8 + col * 2.0)
    yy = Inches(4.1 + row * 0.45)
    chip = add_shape(slide, x, yy, Inches(1.85), Inches(0.35), WHITE, BLUE)
    add_text(slide, c, x + Inches(0.05), yy + Inches(0.02), Inches(1.75), Inches(0.3),
             size=10, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: 기술 스택 & 구현 계획
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, '7. 기술 스택 & 구현 계획', Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         size=28, bold=True, color=DARK)

# Tech stack cards
techs = [
    ('프레임워크', 'Next.js 16 (App Router)\nReact 19 + TypeScript 5', GREEN),
    ('UI', 'Tailwind CSS 4\nshadcn/ui (카드, 입력, 배지)', BLUE),
    ('폼 관리', 'react-hook-form\nzod (유효성 검사 스키마)', AMBER),
    ('상태 관리', 'TanStack Query v5 (서버)\nZustand (인증)', RGBColor(0x8B,0x5C,0xF6)),
]

for i, (t, d, c) in enumerate(techs):
    x = Inches(0.8 + i * 3.1)
    card = add_shape(slide, x, Inches(1.2), Inches(2.8), Inches(1.5), WHITE, RGBColor(0xE2,0xE8,0xF0))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.2), Inches(2.8), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    add_text(slide, t, x + Inches(0.15), Inches(1.23), Inches(2.5), Inches(0.3),
             size=13, bold=True, color=WHITE)
    add_text(slide, d, x + Inches(0.15), Inches(1.7), Inches(2.5), Inches(0.8),
             size=12, color=GRAY)

# Implementation phases
add_text(slide, '구현 단계', Inches(0.8), Inches(3.0), Inches(5), Inches(0.4),
         size=18, bold=True, color=DARK)

phases = [
    ('Phase 1', 'UI 구현', '폼 레이아웃, 칩 버튼, 조건부 필드\n금액 자동 계산, 전화번호 포맷', GREEN),
    ('Phase 2', 'API 연동', 'POST /admin/orders 생성 API 연동\n주소 검색 (Daum API) 연동', BLUE),
    ('Phase 3', '고급 기능', '메시지 파싱 엔진\n최근 발주 불러오기\n경조사어/카드 메시지 검색', AMBER),
    ('Phase 4', '배포 & QA', '빌드/배포 (dev + prod)\n실사용 테스트 & 피드백 반영', RGBColor(0x8B,0x5C,0xF6)),
]

for i, (phase, title, desc, clr) in enumerate(phases):
    yy = Inches(3.5 + i * 0.95)
    
    badge = add_shape(slide, Inches(0.8), yy, Inches(1.2), Inches(0.75), clr)
    add_text(slide, phase, Inches(0.85), yy + Inches(0.05), Inches(1.1), Inches(0.3),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(0.85), yy + Inches(0.35), Inches(1.1), Inches(0.3),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    card = add_shape(slide, Inches(2.2), yy, Inches(10), Inches(0.75), WHITE, RGBColor(0xE2,0xE8,0xF0))
    add_text(slide, desc, Inches(2.4), yy + Inches(0.05), Inches(9.5), Inches(0.65),
             size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN)

add_text(slide, '요약', Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         size=20, color=RGBColor(0xD1,0xFA,0xE5))
add_text(slide, '경쟁사 대비 차별화 포인트', Inches(1), Inches(2.2), Inches(11), Inches(0.8),
         size=36, bold=True, color=WHITE)

summary_items = [
    '카카오톡 메시지 파싱으로 자동 필드 채움 (경쟁사에 없는 고유 기능)',
    '모던 카드 기반 UI — 테이블 레이아웃 탈피',
    '조건부 필드 (근조/축하)로 불필요한 입력 제거',
    'Daum 주소 인라인 검색 (팝업 없이)',
    '실시간 금액 계산 + zod 유효성 검사',
    '완전한 모바일 반응형 (경쟁사 모바일 미지원)',
]

for i, item in enumerate(summary_items):
    yy = Inches(3.3 + i * 0.5)
    add_text(slide, f'✓  {item}', Inches(1.5), yy, Inches(10), Inches(0.45),
             size=16, color=WHITE)

add_text(slide, '달려라 꽃배달 — 더 빠르고 쉬운 주문 등록 경험', 
         Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         size=14, color=RGBColor(0xA7,0xF3,0xD0))

# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
output_path = '/home/blueadm/frontend_web/public/downloads/order-register-proposal.pptx'
prs.save(output_path)
print(f'PPTX saved to {output_path}')
print(f'File size: {os.path.getsize(output_path)} bytes')
